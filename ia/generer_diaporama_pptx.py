# -*- coding: utf-8 -*-
"""
Produit le diaporama de soutenance au format PowerPoint.

Deux consignes se contredisaient et il a fallu trancher. L'ecole demande que le
diaporama n'reprenne pas le plan du memoire mais explique simplement ; le
cahier des charges recu demandait au contraire la fidelite au document. La
lecture retenue est celle-ci : fidelite absolue aux FAITS, aucun chiffre
invente ni arrondi, mais liberte totale sur la STRUCTURE, organisee selon les
questions de l'auditoire et non selon les sections du document.

Meme arbitrage pour les figures. Toutes viennent du memoire, aucune n'est
recreee. Mais seules celles qui restent lisibles en projection sont retenues :
un diagramme de classes ou une matrice de confusion, projetes dix secondes,
ne prouvent rien et coutent une minute. Ils restent dans le memoire, ou le jury
peut les lire.

Le fichier est construit a partir du diaporama existant plutot que de zero, ce
qui conserve la page de garde a l'identique, avec ses logos d'origine et son
gabarit. Recreer une page de garde reviendrait a redessiner des logos, ce que
le cahier des charges interdit a juste titre.

Les transitions et les animations sont ecrites directement dans le XML : la
bibliotheque ne les expose pas. Elles sont donc de vraies animations
PowerPoint, attachees aux formes, et non un effet simule.
"""
import copy
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.oxml import parse_xml

MODELE = Path(r"C:\Users\DELL\Downloads\Memoire_NGuessan_Diby_Soutenance_FINAL (9).pptx")
SORTIE = Path(r"C:\Users\DELL\Downloads\MEMOIRE\Soutenance_SI-ENV.pptx")
FIGURES = Path(r"C:\Users\DELL\AppData\Local\Temp\claude"
               r"\d--etude-soutenance-SI-ENV\3233866b-194c-446a-b8f6-65c65b911c25"
               r"\scratchpad\figures")

BLEU = RGBColor(0x00, 0x4F, 0x9F)
BLEU_SOMBRE = RGBColor(0x00, 0x3A, 0x73)
ORANGE = RGBColor(0xF3, 0x70, 0x21)
ENCRE = RGBColor(0x10, 0x1B, 0x2B)
GRIS = RGBColor(0x55, 0x67, 0x7E)
GRIS_PALE = RGBColor(0x8A, 0x99, 0xAC)
CARTE = RGBColor(0xF2, 0xF5, 0xFA)
FILET = RGBColor(0xD7, 0xDE, 0xE9)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Segoe UI"

MARGE = Inches(1.0)
UTILE = Inches(18.0)
NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ─── Mise en forme ───────────────────────────────────────────────────────

def zone(d, g, h, l, ht):
    tb = d.shapes.add_textbox(g, h, l, ht)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def ecrire(tf, texte, taille, couleur, police=SANS, gras=False,
           interligne=1.2, avant=0, premier=False, align=None):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.line_spacing = interligne
    p.space_before = Pt(avant)
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = texte
    r.font.name = police
    r.font.size = Pt(taille)
    r.font.bold = gras
    r.font.color.rgb = couleur
    return p


def oeil(d, texte):
    tb, tf = zone(d, MARGE, Inches(0.72), UTILE, Inches(0.4))
    ecrire(tf, texte.upper(), 13, ORANGE, gras=True, premier=True)
    return tb


def titre(d, texte, taille=40):
    tb, tf = zone(d, MARGE, Inches(1.28), UTILE, Inches(1.7))
    ecrire(tf, texte, taille, BLEU_SOMBRE, police=SERIF,
           interligne=1.06, premier=True)
    return tb


def chapo(d, texte, haut, taille=17, largeur=Inches(15.0)):
    tb, tf = zone(d, MARGE, haut, largeur, Inches(2.0))
    ecrire(tf, texte, taille, GRIS, interligne=1.45, premier=True)
    return tb


def carte(d, g, h, l, ht, marque, tete, corps=None):
    f = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, g, h, l, ht)
    f.fill.solid()
    f.fill.fore_color.rgb = CARTE
    f.line.color.rgb = FILET
    f.line.width = Pt(1)
    f.shadow.inherit = False
    tf = f.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = tf.margin_bottom = Inches(0.22)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    ecrire(tf, marque.upper(), 11, GRIS_PALE, gras=True, premier=True)
    ecrire(tf, tete, 19, ENCRE, gras=True, interligne=1.15, avant=7)
    if corps:
        ecrire(tf, corps, 14, GRIS, interligne=1.35, avant=7)
    return f


def rangee(d, items, haut, hauteur, colonnes=None):
    n = colonnes or len(items)
    esp = Inches(0.3)
    l = int((UTILE - esp * (n - 1)) / n)
    formes = []
    for i, (m, t, c) in enumerate(items):
        lig, col = divmod(i, n)
        formes.append(carte(d, MARGE + col * (l + esp),
                            haut + lig * (hauteur + esp), l, hauteur, m, t, c))
    return formes


def chiffres(d, items, haut):
    n = len(items)
    esp = Inches(0.4)
    l = int((UTILE - esp * (n - 1)) / n)
    formes = []
    for i, (v, leg) in enumerate(items):
        g = MARGE + i * (l + esp)
        trait = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, g, haut, l, Pt(3))
        trait.fill.solid()
        trait.fill.fore_color.rgb = BLEU
        trait.line.fill.background()
        trait.shadow.inherit = False
        tb, tf = zone(d, g, haut + Inches(0.18), l, Inches(1.9))
        ecrire(tf, v, 52, BLEU_SOMBRE, police=SERIF, interligne=1.0, premier=True)
        ecrire(tf, leg.upper(), 12, GRIS_PALE, gras=True, interligne=1.3, avant=8)
        formes.append(tb)
    return formes


def tableau(d, entetes, lignes, haut, hauteur, parts, retenue=None, regard=False):
    forme = d.shapes.add_table(len(lignes) + 1, len(entetes),
                               MARGE, haut, UTILE, hauteur)
    t = forme.table
    t.first_row = True
    total = sum(parts)
    for i, part in enumerate(parts):
        t.columns[i].width = Emu(int(UTILE * part / total))
    for j, e in enumerate(entetes):
        c = t.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = BLANC
        c.margin_left = c.margin_right = Inches(0.18)
        r = c.text_frame.paragraphs[0].add_run()
        r.text = e.upper()
        r.font.name, r.font.size, r.font.bold = SANS, Pt(12), True
        r.font.color.rgb = ORANGE if (regard and j == len(entetes) - 1) else GRIS_PALE
    for i, ligne in enumerate(lignes, 1):
        for j, v in enumerate(ligne):
            c = t.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = BLANC
            c.margin_left = c.margin_right = Inches(0.18)
            r = c.text_frame.paragraphs[0].add_run()
            r.text = v
            r.font.name, r.font.size = SANS, Pt(16)
            vedette = (retenue is not None and i == retenue)
            if regard and j == len(entetes) - 1:
                r.font.color.rgb, r.font.bold = BLEU, True
            elif regard:
                r.font.color.rgb = ENCRE
            else:
                r.font.color.rgb = ENCRE if vedette else GRIS
                r.font.bold = vedette
    return forme


def figure(d, nom, g, h, hauteur_max, largeur_max):
    """Insere une figure du memoire en respectant ses proportions."""
    chemin = FIGURES / nom
    if not chemin.exists():
        print(f"    figure absente : {nom}")
        return None
    from PIL import Image
    with Image.open(chemin) as im:
        rl, rh = im.size
    ratio = rl / rh
    ht = hauteur_max
    l = int(ht * ratio)
    if l > largeur_max:
        l = largeur_max
        ht = int(l / ratio)
    return d.shapes.add_picture(str(chemin), g + int((largeur_max - l) / 2),
                                h, l, ht)


def numero(d, n, total):
    tb, tf = zone(d, Inches(17.4), Inches(10.4), Inches(1.6), Inches(0.4))
    ecrire(tf, f"{n} / {total}", 12, GRIS_PALE, gras=True, premier=True,
           align=PP_ALIGN.RIGHT)


def notes(d, texte):
    d.notes_slide.notes_text_frame.text = texte


# ─── Transitions et animations, ecrites dans le XML ──────────────────────

AVEC_TRANSITIONS = os.getenv("SANS_TRANSITIONS") != "1"
AVEC_ANIMATIONS = os.getenv("SANS_ANIMATIONS") != "1"


def transition(d, morph=False):
    """Transition native de fondu, greffee dans le XML de la diapositive.

    La bibliotheque n'expose pas les transitions. Le fondu est retenu plutot
    que le morphing : morph n'appartient pas au namespace de 2010 mais a celui
    de 2015, et un element invalide rend le fichier illisible, ce qu'une
    premiere tentative a verifie. Le fondu est de surcroit le seul effet qui se
    comporte pareil sur toutes les versions, ce qui compte quand on presente
    sur une machine que l'on ne connait pas.

    Le noeud se place apres clrMapOvr et avant timing, ordre impose par le
    schema.
    """
    if not AVEC_TRANSITIONS:
        return
    # La page de garde vient du diaporama d'origine et porte deja sa
    # transition. En ajouter une seconde produit deux elements la ou le schema
    # n'en admet qu'un, et PowerPoint refuse alors le fichier entier.
    if d._element.find(qn('p:transition')) is not None:
        return
    xml = f'<p:transition xmlns:p="{NS}" spd="med"><p:fade/></p:transition>'
    noeud = parse_xml(xml)
    timing = d._element.find(qn('p:timing'))
    if timing is not None:
        timing.addprevious(noeud)
    else:
        d._element.append(noeud)


def apparitions(d, formes, duree=500):
    """Fait apparaitre les formes en fondu, l'une apres l'autre, au clic.

    Chaque groupe est un effet d'entree attache a la forme par son
    identifiant. Sans le noeud set qui rend la forme visible, PowerPoint
    afficherait l'element des l'ouverture puis le ferait clignoter.
    """
    if not AVEC_ANIMATIONS or not formes:
        return
    ids = [f.shape_id for f in formes if f is not None]
    if not ids:
        return

    compteur = [3]

    def suivant():
        compteur[0] += 1
        return compteur[0]

    groupes = []
    for sid in ids:
        a, b, c, e, f = (suivant() for _ in range(5))
        groupes.append(f'''
        <p:par><p:cTn id="{a}" fill="hold">
          <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
          <p:childTnLst><p:par><p:cTn id="{b}" fill="hold">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="{c}" presetID="10" presetClass="entr"
                 presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:set><p:cBhvr>
                  <p:cTn id="{e}" dur="1" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                  <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
                <p:animEffect transition="in" filter="fade">
                  <p:cBhvr><p:cTn id="{f}" dur="{duree}"/>
                  <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>
                </p:animEffect>
              </p:childTnLst></p:cTn></p:par></p:childTnLst>
          </p:cTn></p:par></p:childTnLst>
        </p:cTn></p:par>''')

    bld = ''.join(f'<p:bldP spid="{s}" grpId="0"/>' for s in ids)
    xml = f'''<p:timing xmlns:p="{NS}"><p:tnLst>
      <p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst><p:seq concurrent="1" nextAc="seek">
          <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>{''.join(groupes)}</p:childTnLst></p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0">
            <p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0">
            <p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
        </p:seq></p:childTnLst></p:cTn></p:par>
    </p:tnLst><p:bldLst>{bld}</p:bldLst></p:timing>'''
    d._element.append(parse_xml(xml))


# ─── Construction ────────────────────────────────────────────────────────

def construire():
    prs = Presentation(str(MODELE))

    # La page de garde est conservee telle quelle, avec ses logos d'origine.
    # Les trente et une autres sont retirees.
    liste = prs.slides._sldIdLst
    for sid in list(liste)[1:]:
        rid = sid.get(qn('r:id'))
        prs.part.drop_rel(rid)
        liste.remove(sid)
    print(f"page de garde conservee, {len(prs.slides)} diapositive restante")

    vierge = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    TOTAL = 13
    n = [1]

    def nouvelle():
        n[0] += 1
        d = prs.slides.add_slide(vierge)
        for sh in list(d.shapes):
            if sh.has_text_frame and not sh.text_frame.text:
                sh._element.getparent().remove(sh._element)
        fond = d.background.fill
        fond.solid()
        fond.fore_color.rgb = BLANC
        numero(d, n[0], TOTAL)
        return d

    # 2 ─ Le probleme
    d = nouvelle()
    oeil(d, "Le point de départ")
    titre(d, "Une exigence réglementaire tenue sur des fiches papier")
    c = chapo(d, "Le PTUA représente 657,8 milliards de francs et il est classé "
                 "Catégorie 1 par la Banque Africaine de Développement. Ce "
                 "classement impose à l'AGEROUTE un suivi environnemental "
                 "documenté. Dans les faits, il se faisait sur des fiches papier "
                 "et des tableurs.", Inches(3.2))
    f = chiffres(d, [("semaines", "de délai entre le constat et sa remontée"),
                     ("jours", "de ressaisie pour produire un rapport"),
                     ("0", "alerte automatique sur franchissement de seuil")],
                 Inches(6.4))
    transition(d)
    apparitions(d, [c] + f)
    notes(d, "Installer le manque avant de parler de la solution. Les trois "
             "chiffres suffisent, ne pas les commenter un par un. Le jury doit "
             "sentir l'écart entre ce que la BAD exige et ce qui existait.")

    # 3 ─ Objectifs
    d = nouvelle()
    oeil(d, "Ce que je me suis fixé")
    titre(d, "Un objectif général, cinq objectifs spécifiques")
    c = chapo(d, "Concevoir le SI-ENV, un système numérique de suivi "
                 "environnemental des chantiers du PTUA.", Inches(3.2), taille=19)
    f = rangee(d, [
        ("OS 1", "Analyser", "Cartographier les processus de suivi actuels."),
        ("OS 2", "Architecturer", "Définir mobile, serveur, tableau de bord, satellite."),
        ("OS 3", "Développer", "Une application mobile fonctionnant hors ligne."),
        ("OS 4", "Intégrer", "Les données satellitaires via Google Earth Engine."),
        ("OS 5", "Déployer", "Documenter et préparer la mise en service."),
    ], Inches(4.6), Inches(2.6), colonnes=5)
    transition(d)
    apparitions(d, [c] + f)
    notes(d, "Ne pas lire les cinq. Dire que l'objectif général se décline en "
             "cinq étapes, du diagnostic à la mise en production, et enchaîner.")

    # 4 ─ Six manques, six reponses
    d = nouvelle()
    oeil(d, "Ce qui a servi de cahier des charges")
    titre(d, "À chaque manque du terrain, une réponse du système", 38)
    t = tableau(d, ["Limite du dispositif actuel", "Réponse du SI-ENV"],
                [["Détection tardive", "Saisie mobile immédiate sur le terrain"],
                 ["Gravité appréciée sans référentiel", "Classification automatique par intelligence artificielle"],
                 ["Aucune géolocalisation", "Coordonnées GPS et rattachement à l'ouvrage"],
                 ["Données dispersées", "Base spatiale centralisée"],
                 ["Rapport produit à la main", "Génération automatique depuis le tableau de bord"],
                 ["Aucune alerte automatique", "Seuils satellitaires et notifications"]],
                Inches(3.3), Inches(6.4), [42, 58], regard=True)
    transition(d)
    notes(d, "La diapositive qui explique le projet en un écran. Ne pas lire les "
             "douze lignes : dire que chaque limite du terrain a commandé une "
             "réponse, et laisser le tableau parler quelques secondes.")

    # 5 ─ Architecture, avec la figure du memoire
    d = nouvelle()
    oeil(d, "Comment c'est construit")
    titre(d, "Trois clients, un serveur, une base spatiale")
    img = figure(d, "Figure_4-1.png", MARGE, Inches(3.0), Inches(5.6), Inches(11.0))
    f = rangee(d, [
        ("Ce que l'on voit", "Trois applications",
         "Deux mobiles, agent et riverain, et un tableau de bord web."),
        ("Ce qui décide", "Un serveur",
         "Règles métier, habilitations, production des rapports."),
        ("Ce qui conserve", "Une base spatiale",
         "Données, localisation et photographies."),
    ], Inches(3.0), Inches(2.0), colonnes=1)
    for forme in f:
        forme.left = Inches(12.4)
        forme.width = Inches(6.6)
    for i, forme in enumerate(f):
        forme.top = Inches(3.0) + i * Inches(2.3)
    transition(d)
    apparitions(d, f)
    notes(d, "Le point à faire passer : les applications ne parlent jamais "
             "directement à la base, tout passe par le serveur. C'est ce "
             "découplage qui a permis d'ajouter l'application citoyenne en fin "
             "de projet sans toucher aux deux autres.")

    # 6 ─ Les huit profils
    d = nouvelle()
    oeil(d, "Qui s'en sert")
    titre(d, "Huit profils, dont deux qui ne peuvent rien modifier")
    f = rangee(d, [
        ("Terrain", "Responsable Environnement", "Constate et signale."),
        ("Terrain", "Expert HSE", "Traite et valide les actions correctives."),
        ("Pilotage", "Spécialiste Environnement", "Consolide, analyse, produit le rapport."),
        ("Pilotage", "Spécialiste Suivi du PAR", "Instruit les plaintes."),
        ("Consultation seule", "ANDE", "L'agence de tutelle."),
        ("Consultation seule", "BAD", "Le bailleur."),
        ("Riverain", "Plaignant", "Dépose une doléance et la suit."),
        ("Exploitation", "Administrateur", "Comptes et paramétrage."),
    ], Inches(3.3), Inches(2.5), colonnes=4)
    c = chapo(d, "L'agence de tutelle et le bailleur voient les mêmes écrans que "
                 "le spécialiste, sans aucune commande d'écriture. La restriction "
                 "est appliquée par le serveur, pas seulement masquée dans "
                 "l'interface.", Inches(8.9))
    transition(d)
    apparitions(d, [c])
    notes(d, "L'argument de conception le plus fort : celui qui contrôle ne doit "
             "pas pouvoir modifier ce qu'il examine. Insister là-dessus, c'est ce "
             "qu'un jury retient.")

    # 7 ─ Ce qui est livre, avec les captures du memoire
    d = nouvelle()
    oeil(d, "Ce que ça donne")
    titre(d, "Trois applications, un même flux de données")
    a = figure(d, "Figure_5-2.png", MARGE, Inches(3.1), Inches(5.0), Inches(12.0))
    b = figure(d, "Figure_5-2_bis.png", Inches(13.6), Inches(3.1), Inches(5.0), Inches(4.6))
    for forme, libelle in ((a, "L'agent de terrain"), (b, "Le riverain")):
        if forme is None:
            continue
        tb, tf = zone(d, forme.left, Inches(8.35), forme.width, Inches(0.6))
        ecrire(tf, libelle, 19, ENCRE, gras=True, premier=True,
               align=PP_ALIGN.CENTER)
    c = chapo(d, "Saisie géolocalisée en moins d'une minute, qui fonctionne sans "
                 "réseau. Et une seconde application, ouverte aux habitants, dont "
                 "l'accès se conditionne à la présence dans la zone d'influence "
                 "d'un chantier.", Inches(9.1))
    transition(d)
    apparitions(d, [c])
    notes(d, "Si la démonstration en direct est possible, c'est ici qu'elle se "
             "place. Sinon, décrire le parcours : l'agent saisit, ça part au "
             "retour du réseau, le spécialiste le voit. Le canal riverain donne "
             "un support au Mécanisme de Gestion des Plaintes exigé par la BAD.")

    # 8 ─ Le choix du modele, tableau complet du memoire
    d = nouvelle()
    oeil(d, "Première décision de conception")
    titre(d, "La reconnaissance d'images tourne sur le téléphone")
    tableau(d, ["Modèle", "mAP@0.5", "Précision", "Rappel", "F1-Score", "Inférence"],
            [["YOLOv8n", "0,807", "0,797", "0,717", "0,755", "4,3 ms"],
             ["SSD300", "0,612", "0,640", "0,580", "0,608", "185,3 ms"],
             ["Faster R-CNN", "0,685", "0,710", "0,640", "0,673", "312,5 ms"]],
            Inches(3.3), Inches(3.0), [26, 15, 15, 15, 15, 14], retenue=1)
    c = chapo(d, "Le critère décisif n'était pas la précision seule. Faster R-CNN "
                 "n'est pas loin, mais il lui faut soixante-douze fois plus de "
                 "temps. Sur un chantier sans réseau, c'est ce chiffre qui "
                 "tranche.", Inches(7.0))
    transition(d)
    apparitions(d, [c])
    notes(d, "Dire le prix du choix avant son bénéfice. On perd un peu en "
             "précision, on gagne de pouvoir travailler là où le réseau manque. "
             "Ajouter que MobileNetV2 classe la criticité en 8,9 Mo embarqués, "
             "et que le diagnostic reste une aide au classement, jamais une "
             "décision.")

    # 9 ─ Teledetection, avec la figure du memoire
    d = nouvelle()
    oeil(d, "Seconde source d'observation")
    titre(d, "Quatre indices pour voir ce que l'agent ne voit pas")
    figure(d, "Figure_5-9.png", MARGE, Inches(3.1), Inches(5.6), Inches(11.4))
    f = rangee(d, [
        ("Végétation", "NDVI", "Dégradation du couvert autour du chantier."),
        ("Stress hydrique", "NDWI", "Teneur en eau de la végétation."),
        ("Qualité de l'air", "NO2", "Activité des engins et groupes électrogènes."),
        ("Risque terrain", "Pluie et relief", "Eaux stagnantes, donc gîtes à moustiques."),
    ], Inches(3.1), Inches(1.55), colonnes=1)
    for i, forme in enumerate(f):
        forme.left = Inches(12.8)
        forme.width = Inches(6.2)
        forme.top = Inches(3.1) + i * Inches(1.75)
    transition(d)
    apparitions(d, f)
    notes(d, "Rester bref. Le point à faire passer : la donnée satellitaire "
             "complète le terrain, elle ne le remplace pas. Un couvert dégradé "
             "peut signaler un chantier mal maîtrisé comme une saison sèche.")

    # 10 ─ Resultats
    d = nouvelle()
    oeil(d, "Comment je sais que ça marche")
    titre(d, "Ce qui a été mesuré, et non ce qui était espéré")
    f = chiffres(d, [("119", "tests automatisés rejoués à chaque modification"),
                     ("80,7 %", "de précision moyenne sur six classes de déchets"),
                     ("6", "chantiers du PTUA couverts"),
                     ("0", "franc d'hébergement à ce jour")], Inches(4.0))
    # La courbe precision-rappel a ete retiree : reduite a la place disponible,
    # sa legende devenait illisible et elle n'ajoutait rien que les chiffres ne
    # disent mieux. Elle reste dans le memoire, ou elle se lit.
    c = chapo(d, "Le système est déployé et joignable en ligne. Le rapport "
                 "réglementaire, qui demandait plusieurs jours de ressaisie, se "
                 "produit en quelques secondes, et sa remise à l'agence de "
                 "tutelle est tracée.", Inches(7.6), taille=19)
    transition(d)
    apparitions(d, f + [c])
    notes(d, "Trois chiffres suffisent, en annoncer dix revient à n'en faire "
             "retenir aucun. Insister sur le dernier : le système coûte zéro "
             "franc, ce qui était une contrainte du sujet et non un choix par "
             "défaut. Si on demande la courbe précision-rappel, elle est au "
             "mémoire, figure 5.7.")

    # 11 ─ Limites
    d = nouvelle()
    oeil(d, "Ce que ce travail ne démontre pas")
    titre(d, "Quatre limites, dites avant qu'on les relève")
    t = tableau(d, ["", "La limite et sa portée"],
                [["01", "Aucun usage réel. Le système n'a pas été utilisé par des agents qui ne l'ont pas conçu, seul essai qui permettrait de juger de son ergonomie."],
                 ["02", "Corpus d'entraînement. Il vient de sources ouvertes et ne reflète pas les chantiers d'Abidjan, ce qui limite la portée du 80,7 %."],
                 ["03", "Vérification de position. Elle atteste d'une localisation, non d'une résidence."],
                 ["04", "Hébergement. Les offres gratuites conviennent à une validation académique, pas à une mise en service."]],
                Inches(3.3), Inches(5.4), [5, 95])
    transition(d)
    notes(d, "Ralentir ici. Devancer l'objection : un candidat qui énonce ses "
             "limites paraît maîtriser son sujet, celui qui les subit paraît les "
             "avoir manquées. C'est la diapositive qui fait la différence.")

    # 12 ─ Perspectives
    d = nouvelle()
    oeil(d, "La suite")
    titre(d, "Ces limites dessinent le travail qui reste")
    f = rangee(d, [
        ("Répond à la limite 01", "Une phase pilote",
         "Un chantier, les équipes en place, et l'observation de ce qui les bloque réellement."),
        ("Répond à la limite 02", "Un corpus du PTUA",
         "Chaque signalement validé nourrit le jeu d'entraînement, pour un réapprentissage périodique."),
        ("Répond à la limite 04", "Un hébergement de production",
         "Le passage ne demande qu'un changement de paramètres, l'application étant conteneurisée."),
    ], Inches(3.4), Inches(3.2))
    c = chapo(d, "L'architecture ne dépend pas du PTUA. Elle peut équiper "
                 "d'autres projets de l'AGEROUTE par un simple paramétrage du "
                 "référentiel d'ouvrages.", Inches(7.4))
    transition(d)
    apparitions(d, f + [c])
    notes(d, "Chaque piste répond à une limite de la diapositive précédente. Le "
             "lien doit s'entendre, c'est ce qui montre que les limites ont été "
             "pensées et pas seulement subies.")

    # 13 ─ Conclusion
    d = nouvelle()
    oeil(d, "Conclusion")
    titre(d, "Ce que j'en retiens")
    barre = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGE, Inches(3.4), Pt(4), Inches(1.5))
    barre.fill.solid()
    barre.fill.fore_color.rgb = ORANGE
    barre.line.fill.background()
    barre.shadow.inherit = False
    tb, tf = zone(d, MARGE + Inches(0.45), Inches(3.4), Inches(14.0), Inches(1.6))
    ecrire(tf, "Ce projet m'a moins appris à écrire du code qu'à décider.",
           34, ENCRE, police=SERIF, interligne=1.3, premier=True)
    c = chapo(d, "Arbitrer entre un modèle précis et un modèle rapide, accepter "
                 "qu'une fonctionnalité utile ne soit pas prioritaire, "
                 "reconnaître qu'une mesure flatteuse ne prouve rien. Ils m'ont "
                 "aussi appris qu'un outil techniquement irréprochable dont "
                 "personne ne se sert n'a rien résolu.", Inches(5.6))
    tb2, tf2 = zone(d, MARGE, Inches(8.4), Inches(14.0), Inches(0.8))
    ecrire(tf2, "Je vous remercie de votre attention.", 24, BLEU_SOMBRE,
           police=SERIF, premier=True)
    transition(d)
    apparitions(d, [c, tb2])
    notes(d, "Terminer sur le regard, pas sur l'écran. Le silence qui suit vous "
             "appartient : ne le comblez pas.")

    # La page de garde recoit aussi sa transition.
    transition(prs.slides[0])

    prs.save(str(SORTIE))
    print(f"{len(prs.slides)} diapositives")
    print(f"PPTX genere : {SORTIE.name}")


if __name__ == "__main__":
    construire()
